"""Generate SC// proposal / configuration PowerPoint decks.

Decks are derived from the branded SC template at resources/template.pptx so they
inherit its theme (Arial + the SC colour scheme) and slide masters; our content
slides are drawn on the template's BLANK layout. If the template file is missing
(e.g. not deployed), we fall back to a plain blank presentation so exports never
break.
"""

import io
import os
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

from export_gauges import render_util_bars, util_rows, compute_floor_sentence
from recommend import _rec_network_svg
from i18n import translator, font_for

# Brand palette taken from the template theme (resources/template.pptx):
#   dk1 272727 · dk2 113859 · lt2 E9EAF0 · accent1 009ADE · accent2 194F90
#   accent4 3FB748 · accent5 97CAEB · accent6 F78D2C
SC_BLUE = RGBColor(0x00, 0x9A, 0xDE)       # accent1 — primary SC blue
SC_DARK_BLUE = RGBColor(0x11, 0x38, 0x59)  # dk2 — title bar / headings
SC_DEEP_BLUE = RGBColor(0x19, 0x4F, 0x90)  # accent2 — secondary accent / rules
SC_LIGHT_BLUE = RGBColor(0x97, 0xCA, 0xEB)  # accent5 — "SC//" prefix on dark
CHARCOAL = RGBColor(0x27, 0x27, 0x27)      # dk1
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE9, 0xEA, 0xF0)    # lt2
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x3F, 0xB7, 0x48)         # accent4
RED = RGBColor(0xC6, 0x28, 0x28)           # semantic warning (no template red)
ORANGE = RGBColor(0xF7, 0x8D, 0x2C)        # accent6
BORDER_SUBTLE = RGBColor(0xDE, 0xE2, 0xE6)
CARD_BG = RGBColor(0xE9, 0xEA, 0xF0)       # lt2
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)

_TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "resources", "template.pptx")

# Slide-title font — embedded in the template (embeddedFontLst), so it renders
# even where it isn't installed. ExtraLight is a thin weight, so titles aren't bold.
TITLE_FONT = "Martel Sans ExtraLight"


def _new_deck():
    """A fresh deck derived from the SC template (sample slides stripped), or a
    plain blank presentation if the template isn't available."""
    if os.path.exists(_TEMPLATE_PATH):
        prs = Presentation(_TEMPLATE_PATH)
        # Remove the template's sample slides. Drop BOTH the sldId reference and
        # the relationship, otherwise the orphaned slide parts linger and collide
        # with our new slides on save ("Duplicate name: slide1.xml").
        sld_id_lst = prs.slides._sldIdLst
        for sld_id in list(sld_id_lst):
            rid = sld_id.get(qn("r:id"))
            if rid:
                prs.part.drop_rel(rid)
            sld_id_lst.remove(sld_id)
    else:
        prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _blank_layout(prs):
    """The template's BLANK layout (falls back to the standard blank)."""
    for layout in prs.slide_layouts:
        if (layout.name or "").strip().upper() == "BLANK":
            return layout
    # python-pptx default template: layout 6 is "Blank".
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]


def _content_layout(prs):
    """The branded light content layout ('b. blank light wide'). Its background
    is a clean light gradient with only a small bottom-right corner + // logo (no
    frame, no right-edge accent), so our tables clear the branding. We draw on top
    of it rather than painting our own background. Searches all masters (this
    layout lives on the second master) and falls back to BLANK if absent."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if (layout.name or "").strip().lower() == "b. blank light wide":
                return layout
    return _blank_layout(prs)


def _svg_to_png_bytes(svg, out_width=2000):
    """Rasterise an SVG string to PNG bytes. Production uses cairosvg; on the dev
    mac (no cairo) it falls back to qlmanage so the export can be tested locally.
    Returns None if neither is available (the caller then skips the diagram)."""
    try:
        import cairosvg
        return cairosvg.svg2png(bytestring=svg.encode("utf-8"), output_width=out_width)
    except Exception:
        pass
    try:  # macOS dev fallback
        import subprocess, tempfile, os
        with tempfile.TemporaryDirectory() as d:
            sp = os.path.join(d, "d.svg")
            with open(sp, "w") as f:
                f.write(svg)
            subprocess.run(["qlmanage", "-t", "-s", str(out_width), "-o", d, sp],
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=30)
            png = sp + ".png"
            return open(png, "rb").read() if os.path.exists(png) else None
    except Exception:
        return None


def _slide_network(prs, recommendation, t, lang="en"):
    """A 'Cluster Network' slide with the recommendation's diagram, scaled to fit
    the content area (preserving the SVG's aspect ratio). Skipped if there's no
    diagram or no rasteriser available."""
    # Regenerate the diagram in the document language (the stored network_svg was
    # rendered in the sizing-time default); fall back to the stored one if needed.
    svg = _rec_network_svg(recommendation, lang) or recommendation.get("network_svg")
    if not svg:
        return
    png = _svg_to_png_bytes(svg)
    if not png:
        return
    slide = _add_slide(prs)
    _add_title(slide, t("export.pptx.cluster_network"), recommendation.get("model", ""),
               lang=lang)
    m = re.search(r'viewBox="0 0 ([\d.]+) ([\d.]+)"', svg)
    vw, vh = (float(m.group(1)), float(m.group(2))) if m else (1200.0, 800.0)
    avail_l, avail_t, avail_w, avail_h = 0.6, 1.7, 12.1, 5.3
    scale = min(avail_w / vw, avail_h / vh)
    bw, bh = vw * scale, vh * scale
    left = avail_l + (avail_w - bw) / 2
    top = avail_t + (avail_h - bh) / 2
    slide.shapes.add_picture(io.BytesIO(png), Inches(left), Inches(top),
                             Inches(bw), Inches(bh))


def generate_proposal(summary, recommendation, projection, source_perf=None, lang="en"):
    t = translator(lang)
    prs = _new_deck()

    _slide_current_env(prs, summary, t, lang)
    _slide_workload(prs, summary, t, lang)
    _slide_proposal(prs, recommendation, projection, t, lang)
    _slide_sizing(prs, recommendation, summary, t, lang)
    _slide_benchmark(prs, recommendation, source_perf, t, lang)
    _slide_network(prs, recommendation, t, lang)
    _slide_projection(prs, summary, recommendation, projection, t, lang)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_config_slide(result, lang="en"):
    t = translator(lang)
    prs = _new_deck()

    slide = _add_slide(prs)
    mode = result.get("mode", "appliance")
    node_count = result["node_count"]
    so = result.get("storage_only")
    if so:
        nodes_label = t("export.pptx.nodes_hci_storage_only",
                        hci=node_count, so=so['count'],
                        total=result.get('total_node_count', node_count))
    else:
        nodes_label = t("export.pptx.nodes_count", count=node_count)
    num_cl = result.get("num_clusters", 1)
    if num_cl > 1:
        layout = result.get("cluster_layout", [])
        nodes_label += "  —  " + t("export.pptx.clusters_layout",
                                   count=num_cl, layout=' + '.join(map(str, layout)))
    pn = result["per_node"]
    cl = result["cluster_total"]
    n1 = result["n_minus_1"]

    if mode == "appliance":
        title = t("export.pptx.configuration_model", model=result.get('model', ''))
        subtitle_parts = [nodes_label]
        if result.get("form_factor"):
            subtitle_parts.append(result["form_factor"])
        if result.get("chassis"):
            subtitle_parts.append(result["chassis"])
        _add_title(slide, title, "  —  ".join(subtitle_parts), lang=lang)

        node_rows = [
            [t("export.common.per_node"), ""],
            ["CPU", pn.get("cpu", "")],
            [t("export.common.cores"), str(pn["cores"])],
            [t("export.common.threads"), str(pn["threads"])],
            [t("export.common.clock_speed"), f"{pn['ghz']} GHz"],
            [t("export.common.ram"), _fmt_ram(pn["ram_gb"])],
            [t("export.common.raw_storage"), f"{pn['raw_storage_tb']} TB"],
        ]
    else:
        title = t("export.pptx.configuration_software_only")
        storage_type = result.get("storage_type", "")
        _add_title(slide, title,
                   t("export.pptx.software_only_subtitle",
                     nodes=nodes_label, storage_type=storage_type,
                     drives=pn.get('disk_count', 0)),
                   lang=lang)

        node_rows = [
            [t("export.common.per_node"), ""],
            [t("export.common.cores"), str(pn["cores"])],
            [t("export.common.threads"), str(pn["threads"])],
            [t("export.common.clock_speed"), f"{pn['ghz']} GHz"],
            [t("export.common.ram"), _fmt_ram(pn["ram_gb"])],
            [t("export.pptx.drives"), str(pn.get("disk_count", 0))],
            [t("export.common.raw_storage"), f"{pn['raw_storage_tb']} TB"],
        ]

    _add_table(slide, 0.6, 1.6, 4.0, node_rows, [1.5, 2.5])

    total_rows = [
        [t("export.common.cluster_total"), ""],
        [t("export.common.cores"), str(cl["cores"])],
        [t("export.common.threads"), str(cl["threads"])],
        ["GHz", str(cl["total_ghz"])],
        [t("export.common.ram"), _fmt_ram(cl["ram_gb"])],
        [t("export.common.raw_storage"), f"{cl['raw_storage_tb']} TB"],
        [t("export.common.usable_storage"), f"{cl['usable_storage_tb']} TB"],
    ]
    _add_table(slide, 4.8, 1.6, 4.0, total_rows, [1.5, 2.5])

    if result.get("single_node"):
        # No peer node to fail over to — N-1 is meaningless. Replace the figures
        # with a greyed-out no-redundancy notice.
        _add_no_redundancy_box(slide, 9.0, 1.6, 4.0,
                               t("export.common.no_redundancy"), t)
    else:
        n1_rows = [
            [t("export.common.n1_available"), ""],
            [t("export.common.cores"), str(n1["cores"])],
            [t("export.common.threads"), str(n1["threads"])],
            ["GHz", str(n1["total_ghz"])],
            [t("export.common.ram"), _fmt_ram(n1["ram_gb"])],
            [t("export.common.usable_storage"), f"{n1['usable_storage_tb']} TB"],
        ]
        _add_table(slide, 9.0, 1.6, 4.0, n1_rows, [1.5, 2.5])

    if so:
        _add_storage_only_note(slide, so, 4.7, t)

    # Append the cluster network diagram as its own slide (manual builder).
    _slide_network(prs, result, t, lang)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _add_slide(prs):
    # Use the branded content layout and let its background show — no white fill.
    return prs.slides.add_slide(_content_layout(prs))


def _add_title(slide, text, subtitle=None, lang="en"):
    # No bar, no "SC//" prefix — the template's branded background (corner + //
    # logo) carries the SC mark. Title sits lower to match the template's title
    # height, in Martel Sans ExtraLight; subtitle beneath it.
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12.2), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(26)
    run.font.bold = False
    run.font.name = font_for(lang, TITLE_FONT)
    run.font.color.rgb = SC_DARK_BLUE

    # Thin dark-blue rule from just after the title to the right edge of the slide.
    # y is the title's CAP vertical centre, computed from Martel Sans ExtraLight's
    # real metrics (ascent 1.15em, capHeight 0.68em) so it lands the same in both
    # PowerPoint and the LibreOffice PDF (both now use the installed font, whose
    # hhea/typo/win metrics all agree). For a 26pt title in a box at top=0.7":
    #   baseline = 0.7 + 0.05 inset + 0.415 ascent ≈ 1.165"
    #   cap centre = baseline − capHeight/2 (0.123") ≈ 1.04"  (old 1.01 sat too high)
    title_end = 0.6 + len(text) * 0.18
    line_x1 = min(title_end + 0.5, 12.0)
    line_y = 1.04
    if line_x1 < 13.0:
        rule = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(line_x1), Inches(line_y),
            Inches(13.333), Inches(line_y))
        rule.line.color.rgb = SC_DARK_BLUE
        rule.line.width = Pt(1)
        rule.shadow.inherit = False  # the template's connector style adds a shadow

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.18), Inches(12.2), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(13)
        run2.font.color.rgb = MID_GRAY


def _add_card(slide, left, top, width, height, label, value, accent=False, lang="en"):
    # For CJK languages the theme font (Arial) can't render the translated card
    # label, so switch those runs to the CJK-capable font; Latin scripts keep the
    # theme default (font_for returns the given latin name unchanged).
    cjk_font = font_for(lang, None)
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8)
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)

    p_label = tf.paragraphs[0]
    run_l = p_label.add_run()
    run_l.text = label
    run_l.font.size = Pt(10)
    run_l.font.color.rgb = MID_GRAY
    if cjk_font:
        run_l.font.name = cjk_font

    p_val = tf.add_paragraph()
    run_v = p_val.add_run()
    run_v.text = str(value)
    run_v.font.size = Pt(18)
    run_v.font.bold = True
    run_v.font.color.rgb = SC_BLUE if accent else CHARCOAL
    if cjk_font:
        run_v.font.name = cjk_font


def _add_no_redundancy_box(slide, left, top, width, msg, t):
    """Greyed-out N-1 replacement for a Single Node System: a header and the
    no-redundancy notice, styled to read as a warning rather than data."""
    height = 2.0
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = Pt(10)
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)

    p_label = tf.paragraphs[0]
    run_l = p_label.add_run()
    run_l.text = t("export.pptx.n1_update_failure")
    run_l.font.size = Pt(11)
    run_l.font.bold = True
    run_l.font.color.rgb = MID_GRAY

    p_head = tf.add_paragraph()
    p_head.space_before = Pt(8)
    run_h = p_head.add_run()
    run_h.text = t("export.common.no_redundancy_heading")
    run_h.font.size = Pt(14)
    run_h.font.bold = True
    run_h.font.color.rgb = RED

    p_msg = tf.add_paragraph()
    p_msg.space_before = Pt(6)
    run_m = p_msg.add_run()
    run_m.text = msg
    run_m.font.size = Pt(10.5)
    run_m.font.color.rgb = CHARCOAL


HEADER_BG = RGBColor(0xE9, 0xEA, 0xF0)   # lt2 — table header band
ROW_ALT_BG = RGBColor(0xF6, 0xF7, 0xFA)  # subtle zebra stripe


def _set_cell_border(tc_pr, side, width_pt, color_hex):
    from pptx.oxml.ns import qn
    tag = {"bottom": "a:lnB", "top": "a:lnT", "left": "a:lnL", "right": "a:lnR"}[side]
    for old in tc_pr.findall(qn(tag)):
        tc_pr.remove(old)
    ln = tc_pr.makeelement(qn(tag), {})
    if color_hex is None:
        ln.append(ln.makeelement(qn("a:noFill"), {}))
    else:
        ln.set("w", str(Pt(width_pt)))
        sf = ln.makeelement(qn("a:solidFill"), {})
        sf.append(sf.makeelement(qn("a:srgbClr"), {"val": color_hex}))
        ln.append(sf)
    tc_pr.append(ln)


def _add_table(slide, left, top, width, rows_data, col_widths=None):
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows_data else 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(0.38 * rows))
    table = table_shape.table

    tbl_pr = table._tbl.tblPr
    for attr in ["bandRow", "bandCol", "firstRow", "lastRow"]:
        tbl_pr.attrib[attr] = "0"

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.margin_top = Pt(5)
            cell.margin_bottom = Pt(5)
            cell.margin_left = Pt(8)

            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = HEADER_BG
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = ROW_ALT_BG
            else:
                cell.fill.fore_color.rgb = WHITE

            tc_pr = cell._tc.get_or_add_tcPr()
            _set_cell_border(tc_pr, "left", None, None)
            _set_cell_border(tc_pr, "right", None, None)
            _set_cell_border(tc_pr, "top", None, None)

            if r == 0:
                _set_cell_border(tc_pr, "bottom", 1.5, "009ADE")
            elif r < rows - 1:
                _set_cell_border(tc_pr, "bottom", 0.5, "DEE2E6")
            else:
                _set_cell_border(tc_pr, "bottom", None, None)

            for paragraph in cell.text_frame.paragraphs:
                if r == 0:
                    paragraph.font.size = Pt(11)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = SC_DARK_BLUE
                elif c == 0:
                    paragraph.font.size = Pt(11)
                    paragraph.font.color.rgb = MID_GRAY
                else:
                    paragraph.font.size = Pt(11)
                    paragraph.font.color.rgb = CHARCOAL

    return table_shape


def _fmt_ram(gb):
    if gb >= 1024:
        return f"{gb / 1024:.1f} TB"
    return f"{gb} GB"


def _storage_only_desc(so):
    """One-line CPU/RAM descriptor for a storage-only-node block. Handles both
    the appliance/recommendation shape (has 'cpu') and the validated shape
    (cores/threads/ghz)."""
    cpu = so.get("cpu")
    if not cpu:
        cpu = f"{so.get('cores', 0)} cores @ {so.get('ghz', 0)}GHz"
    # The block's CPU desc carries a "1 x" prefix (each storage-only node is
    # single-socket); drop it so it doesn't read "2 × 1 x ..." after the count.
    for prefix in ("1 x ", "1 × "):
        if cpu.startswith(prefix):
            cpu = cpu[len(prefix):]
            break
    return cpu


def _add_storage_only_note(slide, so, top, t):
    """Render the storage-only-node summary line. ``so`` is the storage_only
    block; ``top`` is the vertical position in inches."""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = t("export.pptx.storage_only_label") + " "
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = SC_BLUE
    run2 = p.add_run()
    run2.text = t(
        "export.pptx.storage_only_desc",
        count=so['count'],
        cpu=_storage_only_desc(so),
        ram=_fmt_ram(so['ram_gb']),
        raw=so['raw_storage_tb'],
    )
    run2.font.size = Pt(11)
    run2.font.color.rgb = CHARCOAL


def _fmt_num(n):
    if isinstance(n, float):
        return f"{n:,.1f}"
    return f"{n:,}"


# ── Slide 1: Current Environment ─────────────────────────────────────────────

def _slide_current_env(prs, s, t, lang="en"):
    slide = _add_slide(prs)
    _add_title(slide, t("export.pptx.current_environment"),
               f"{s.get('current_platform', '')}  —  {s.get('cluster_name', '')}",
               lang=lang)

    y = 2.35
    cards = [
        (t("export.pptx.hosts"), s.get("host_count", 0)),
        (t("export.pptx.total_cores"), _fmt_num(s.get("total_host_cores", 0))),
        (t("export.pptx.total_threads"), _fmt_num(s.get("total_host_threads", 0))),
        (t("export.pptx.total_ghz"), _fmt_num(s.get("total_host_ghz", 0))),
        (t("export.pptx.total_ram"), _fmt_ram(s.get("total_host_ram_gb", 0))),
    ]
    for i, (label, val) in enumerate(cards):
        _add_card(slide, 0.6 + i * 2.5, y, 2.3, 0.9, label, val, lang=lang)

    y2 = 3.55
    perf = [
        (t("export.pptx.peak_cpu_pct"), f"{s.get('peak_cpu_pct', 0)}%"),
        (t("export.pptx.avg_cpu_pct"), f"{s.get('avg_cpu_pct', 0)}%"),
        (t("export.pptx.peak_memory_pct"), f"{s.get('peak_mem_pct', 0)}%"),
        (t("export.pptx.avg_memory_pct"), f"{s.get('avg_mem_pct', 0)}%"),
    ]
    for i, (label, val) in enumerate(perf):
        _add_card(slide, 0.6 + i * 2.5, y2, 2.3, 0.9, label, val, lang=lang)

    y3 = 4.75
    iops = [
        (t("export.pptx.avg_iops"), _fmt_num(s.get("total_avg_iops", 0))),
        (t("export.pptx.peak_iops"), _fmt_num(s.get("total_peak_iops", 0))),
    ]
    p95 = s.get("p95_iops", 0)
    if p95 and p95 > 0:
        iops.append((t("export.pptx.p95_iops"), _fmt_num(p95)))
    iops.append((t("export.pptx.nic_speed"), f"{s.get('nic_speed_mbps', 0) / 1000:.0f} GbE"))

    for i, (label, val) in enumerate(iops):
        _add_card(slide, 0.6 + i * 2.5, y3, 2.3, 0.9, label, val, lang=lang)

    y4 = 5.95
    ratio = s.get("vcpu_per_core_ratio", 0)
    if ratio > 0:
        _add_card(slide, 0.6, y4, 2.3, 0.9, t("export.pptx.vcpu_core_ratio"),
                  f"{ratio:.2f} : 1", accent=True, lang=lang)


# ── Slide 2: Workload Consumption ────────────────────────────────────────────

def _slide_workload(prs, s, t, lang="en"):
    slide = _add_slide(prs)
    _add_title(slide, t("export.pptx.workload_analysis"),
               t("export.pptx.workload_subtitle",
                 active=s.get('active_vms', 0), total=s.get('total_vms', 0)),
               lang=lang)

    y = 2.35
    compute = [
        (t("export.pptx.total_vcpus"), _fmt_num(s.get("total_vcpus", 0)), True),
        (t("export.pptx.provisioned_ram"), _fmt_ram(s.get("total_vm_provisioned_memory_gb", 0)), True),
        (t("export.pptx.used_ram"), _fmt_ram(s.get("total_vm_used_memory_gb", 0)), False),
    ]
    for i, (label, val, accent) in enumerate(compute):
        _add_card(slide, 0.6 + i * 4.0, y, 3.7, 1.0, label, val, accent, lang=lang)

    y2 = 3.65
    storage = [
        (t("export.pptx.provisioned_storage"), f"{s.get('total_vm_provisioned_storage_tb', 0)} TiB", False),
        (t("export.pptx.datastore_used"), f"{s.get('datastore_used_tb', 0)} TiB", True),
        (t("export.pptx.datastore_total"), f"{s.get('datastore_total_tb', 0)} TiB", False),
    ]
    for i, (label, val, accent) in enumerate(storage):
        _add_card(slide, 0.6 + i * 4.0, y2, 3.7, 1.0, label, val, accent, lang=lang)

    y3 = 5.15
    ratio = s.get("vcpu_per_core_ratio", 0)
    if ratio > 0:
        _add_card(slide, 0.6, y3, 5.5, 1.0,
                  t("export.pptx.current_virtualization_ratio"),
                  t("export.pptx.virtualization_ratio_value",
                    ratio=f"{ratio:.2f}", vcpus=s.get('total_vcpus', 0),
                    cores=s.get('total_host_cores', 0)),
                  accent=True, lang=lang)


# ── Slide 3: Proposed Configuration ──────────────────────────────────────────

def _slide_proposal(prs, r, projection=None, t=None, lang="en"):
    if t is None:
        t = translator(lang)
    slide = _add_slide(prs)

    num_cl = r.get("num_clusters", 1)
    layout = r.get("cluster_layout", [r["node_count"]])
    layout_str = " + ".join(str(x) for x in layout)
    if num_cl > 1:
        cluster_desc = t("export.pptx.cluster_desc_multi", count=num_cl, layout=layout_str)
    else:
        cluster_desc = t("export.pptx.cluster_desc_single")

    if r.get("validated_only"):
        model_label = r["model"]
    elif r.get("validated"):
        model_label = t("export.pptx.validated_based_off", model=r['model'])
    else:
        model_label = r["model"]
    so = r.get("storage_only")
    if so:
        nodes_label = t("export.pptx.nodes_hci_plus_so",
                        hci=r.get('hci_node_count', r['node_count']), so=so['count'])
    else:
        nodes_label = t("export.pptx.nodes_plain", count=r['node_count'])
    _add_title(slide, t("export.pptx.proposed_model", model=model_label),
               f"{nodes_label}  —  {cluster_desc}  —  {r['form_factor']}  —  {r['chassis']}",
               lang=lang)

    iops = r.get("iops") or {}

    node_rows = [
        [t("export.pptx.per_hci_node") if so else t("export.common.per_node"), ""],
        ["CPU", r["cpu"]],
        [t("export.common.cores"), str(r["cores_per_node"])],
        [t("export.common.threads"), str(r["threads_per_node"])],
        [t("export.common.ram"), _fmt_ram(r["ram_per_node_gb"])],
        [t("export.common.storage"), r["storage_config"]["desc"]],
    ]
    if iops:
        node_rows.append([t("export.pptx.net_iops"), f"{iops['per_node']:,}"])
    _add_table(slide, 0.6, 2.35, 4.0, node_rows, [1.5, 2.5])

    tot = r["totals"]
    total_rows = [
        [t("export.common.cluster_total"), ""],
        [t("export.common.cores"), str(tot["cores"])],
        [t("export.common.threads"), str(tot["threads"])],
        ["GHz", str(tot["total_ghz"])],
        [t("export.common.ram"), _fmt_ram(tot["ram_gb"])],
        [t("export.common.raw_storage"), f"{tot['raw_storage_tb']} TB"],
        [t("export.common.usable_storage"), f"{tot['usable_storage_tb']} TB"],
    ]
    if iops:
        total_rows.append([t("export.pptx.net_iops"), f"{iops['total']:,}"])
    _add_table(slide, 4.8, 2.35, 4.0, total_rows, [1.5, 2.5])

    n = r["n_minus_1"]
    if num_cl > 1:
        n1_label = t("export.pptx.n1_spares", count=num_cl)
    else:
        n1_label = t("export.common.n1_available")
    n1_rows = [
        [n1_label, ""],
        [t("export.common.cores"), str(n["cores"])],
        [t("export.common.threads"), str(n["threads"])],
        ["GHz", str(n["total_ghz"])],
        [t("export.common.ram"), _fmt_ram(n["ram_gb"])],
        [t("export.common.usable_storage"), f"{n['usable_storage_tb']} TB"],
    ]
    if iops:
        n1_rows.append([t("export.pptx.net_iops"), f"{iops['n_minus_1']:,}"])
    _add_table(slide, 9.0, 2.35, 4.0, n1_rows, [1.5, 2.5])

    if so:
        _add_storage_only_note(slide, so, 5.25, t)

    _add_card(slide, 0.6, 5.5, 3.5, 0.9,
              t("export.pptx.vcpu_core_ratio_n1"),
              f"{r['vcpu_ratio']:.2f} : 1", accent=True, lang=lang)

    # Net IOPS headroom vs the workload's measured demand (informational).
    demand = (projection or {}).get("iops_demand") or {}
    if iops and (demand.get("p95") or demand.get("avg")):
        metric = "P95" if demand.get("p95") else "Avg"
        value = demand.get("p95") or demand.get("avg")
        ratio = iops["total"] / value if value else 0
        _add_card(slide, 4.3, 5.5, 4.0, 0.9,
                  t("export.pptx.net_iops_headroom_vs", metric=metric),
                  f"{ratio:.1f}x  ({iops['total']:,} net / {value:,} demand)", lang=lang)

    # Derivation footnote — the PPTX is the one place we show how net IOPS is
    # reached (raw drive IOPS, SCRIBE derating, RF write-amplification).
    if iops and iops.get("raw_per_node"):
        note = t("export.pptx.net_iops_derivation",
                 raw=f"{iops['raw_per_node']:,}",
                 derating=f"{iops['derating_pct']:.0f}",
                 derated=f"{iops['derated_per_node']:,}",
                 write_amp=iops['write_amp'],
                 per_node=f"{iops['per_node']:,}",
                 nodes=r['node_count'])
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(11.2), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = note
        run.font.size = Pt(9)
        run.font.color.rgb = MID_GRAY


def _add_textbox(slide, left, top, width, height, lines, lang="en"):
    """lines: list of (text, size_pt, color, bold). Each becomes a paragraph."""
    cjk_font = font_for(lang, None)
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        if cjk_font:
            run.font.name = cjk_font
        p.space_after = Pt(4)
    return box


# ── Sizing rationale: utilization bars + how the node count was reached ──────

def _slide_sizing(prs, r, s, t=None, lang="en"):
    if t is None:
        t = translator(lang)
    u = r.get("utilization")
    if not u:
        return
    slide = _add_slide(prs)
    _add_title(slide, t("export.pptx.sizing_rationale"), r.get("model", ""), lang=lang)

    det = r.get("determinant") or {}
    binding = det.get("resource", "")
    rows, any_ha = util_rows(u)

    top = 2.7   # bars + derivation sit lower, clear of the title rule
    if rows:
        png = render_util_bars(rows, limiting_key=binding, any_ha=any_ha, lang=lang)
        iw, ih = Image.open(io.BytesIO(png)).size
        disp_w = 12.1
        disp_h = disp_w * ih / iw
        slide.shapes.add_picture(io.BytesIO(png), Inches(0.6), Inches(top),
                                 Inches(disp_w), Inches(disp_h))
        top += disp_h + 0.35

    lines = [(t("export.pptx.how_this_was_sized"), 14, SC_DARK_BLUE, True)]
    res = det.get("resource")
    hr = det.get("headroom_pct")
    if res == "CPU":
        ratio = r.get("vcpu_ratio") or 0
        vcpus = s.get("total_vcpus")
        lines.append((t("export.pptx.determined_by_cpu",
                        vcpus=vcpus, ratio=f"{ratio:.2f}",
                        required=f"{det.get('required'):.0f}",
                        achieved=f"{det.get('achieved'):.0f}",
                        headroom=f"{hr:.1f}"), 12, CHARCOAL, False))
    elif res in ("RAM", "Storage"):
        unit = det.get("unit", "")
        res_name = t("export.common.ram") if res == "RAM" else t("export.common.storage")
        lines.append((t("export.pptx.determined_by_resource",
                        resource=res_name, required=det.get('required'), unit=unit,
                        achieved=det.get('achieved'),
                        headroom=f"{hr:.1f}"), 12, CHARCOAL, False))
    elif res == "Compute":
        cf = r.get("compute_floor") or {}
        util = cf.get("source_cpu_util_pct", 100)
        lines.append((t("export.pptx.determined_by_compute",
                        achieved=f"{det.get('achieved'):.0f}",
                        util=f"{util:.0f}"), 12, CHARCOAL, False))
    else:
        lines.append((t("export.pptx.sized_to_minimum"),
                      12, CHARCOAL, False))
    # Show the compute-floor coverage even when another resource was binding, so
    # the deck reflects that sizing is performance-aware.
    if res != "Compute":
        cfs = compute_floor_sentence(r, lang)
        if cfs:
            lines.append((cfs, 11, MID_GRAY, False))
    lines.append((t("export.pptx.bar_legend"), 11, MID_GRAY, False))
    _add_textbox(slide, 0.6, top, 12.1, 1.8, lines, lang=lang)


# ── Benchmark performance: current environment vs recommended cluster ────────

def _slide_benchmark(prs, r, source_perf, t=None, lang="en"):
    if t is None:
        t = translator(lang)
    tgt = (r.get("totals") or {}).get("perf_index")
    if not source_perf or not source_perf.get("total_specrate") or not tgt:
        return
    slide = _add_slide(prs)
    _add_title(slide, t("export.pptx.performance_vs_current"), r.get("model", ""), lang=lang)

    src_total = source_perf["total_specrate"]
    ratio = tgt / src_total if src_total else 0

    _add_textbox(slide, 0.6, 2.3, 7.0, 0.4,
                 [(t("export.pptx.where_you_are_now"), 13, SC_DARK_BLUE, True)], lang=lang)
    src_rows = [[t("export.pptx.your_current_cpus"), t("export.pptx.sockets"),
                 t("export.pptx.score"), "SPECrate"]]
    used_pm = False
    for c in source_perf.get("cpus", []):
        is_pm = c.get("type") == "passmark"
        used_pm = used_pm or is_pm
        score_lbl = f"{_fmt_num(c['score'])} {'PassMark' if is_pm else 'SPECrate'}"
        src_rows.append([c.get("model", ""), str(c.get("sockets", "")),
                         score_lbl, _fmt_num(c.get("total", 0))])
    src_rows.append([t("export.pptx.total_environment"), "", "", _fmt_num(src_total)])
    _add_table(slide, 0.6, 2.75, 7.0, src_rows, [3.4, 1.0, 1.6, 1.0])

    _add_textbox(slide, 8.0, 2.3, 4.7, 0.4,
                 [(t("export.pptx.where_youre_going"), 13, SC_DARK_BLUE, True)], lang=lang)
    used_pm = used_pm or bool(r.get("cpu_perf_is_passmark"))
    tgt_rows = [[t("export.pptx.recommended_cluster"), ""],
                [t("export.pptx.cpu_per_node"), r.get("cpu", "")],
                [t("export.pptx.nodes"), str(r.get("node_count", ""))],
                ["Cluster SPECrate2017", _fmt_num(tgt)]]
    _add_table(slide, 8.0, 2.75, 4.7, tgt_rows, [2.3, 2.4])

    if ratio >= 1:
        verdict = t("export.pptx.verdict_multiple", ratio=f"{ratio:.1f}")
    else:
        verdict = t("export.pptx.verdict_percent", pct=round(ratio * 100))
    _add_card(slide, 0.6, 5.15, 6.2, 0.95, t("export.pptx.in_a_benchmark_deliver"),
              verdict, accent=True, lang=lang)

    foot = [(t("export.pptx.specrate_explainer"), 11, MID_GRAY, False)]
    if used_pm:
        foot.append((t("export.pptx.passmark_note"), 10, MID_GRAY, False))
    foot.append((t("export.pptx.benchmark_disclaimer"), 10, MID_GRAY, False))
    _add_textbox(slide, 0.6, 6.25, 12.1, 1.2, foot, lang=lang)


# ── Slide 4: Growth Projection ───────────────────────────────────────────────

def _slide_projection(prs, s, r, p, t=None, lang="en"):
    if t is None:
        t = translator(lang)
    slide = _add_slide(prs)
    full_cluster = r.get("sized_full_cluster", False)
    cpu_basis = t("export.pptx.cpu_basis_full") if full_cluster else "N-1"
    _add_title(slide, t("export.pptx.capacity_planning", years=p['years']),
               t("export.pptx.projection_subtitle",
                 growth=p['growth_pct'], snapshot=p['snapshot_pct'],
                 factor=p['growth_factor'], cpu_basis=cpu_basis),
               lang=lang)

    n1 = r["n_minus_1"]

    headers = [t("export.common.resource"), t("export.common.current"),
               t("export.pptx.year_projected", years=p['years']),
               t("export.common.proposed_n1"), t("export.common.headroom")]

    # vCPU is sized against the full cluster when that mode is on; RAM, GHz and
    # storage remain N-1, so only this row's basis changes.
    cpu_basis_cores = r["totals"]["cores"] if full_cluster else n1["cores"]
    vcpu_headroom = cpu_basis_cores * r["vcpu_ratio"] - p["projected_vcpus"]
    ram_headroom = n1["ram_gb"] - p["projected_ram_gb"]
    stor_headroom = n1["usable_storage_tb"] - p["projected_storage_tb"]
    ghz_headroom = n1["total_ghz"] - p.get("projected_ghz", 0)

    rows = [
        headers,
        ["vCPUs", _fmt_num(p["base_vcpus"]), _fmt_num(p["projected_vcpus"]),
         t("export.pptx.vcpu_basis", cores=_fmt_num(cpu_basis_cores),
           ratio=f"{r['vcpu_ratio']:.1f}")
         + (" (N)" if full_cluster else ""),
         t("export.pptx.vcpu_headroom", value=_fmt_num(int(vcpu_headroom)))],
        [t("export.common.ram"), _fmt_ram(p["base_ram_gb"]), _fmt_ram(p["projected_ram_gb"]),
         _fmt_ram(n1["ram_gb"]),
         f"+{_fmt_ram(round(ram_headroom, 1))}"],
        ["GHz", f"{p.get('base_ghz', 0)}", f"{p.get('projected_ghz', 0)}",
         f"{n1['total_ghz']}",
         f"+{round(ghz_headroom, 1)} GHz"],
        [t("export.common.storage"), f"{p['base_storage_tb']} TB",
         t("export.pptx.storage_projected_incl_snapshots", value=p['projected_storage_tb']),
         t("export.pptx.storage_usable", value=n1['usable_storage_tb']),
         f"+{round(stor_headroom, 2)} TB"],
    ]

    table_shape = _add_table(slide, 0.6, 2.6, 12.1, rows,
                              [1.5, 2.2, 3.0, 3.0, 2.4])

    y = 4.9
    params = [
        (t("export.pptx.growth_rate"), t("export.pptx.growth_rate_value", pct=p['growth_pct'])),
        (t("export.pptx.growth_factor"),
         t("export.pptx.growth_factor_value", factor=p['growth_factor'], years=p['years'])),
        (t("export.pptx.snapshot_overhead"),
         t("export.pptx.snapshot_overhead_value", base=p['snapshot_pct'],
           target=p.get('snapshot_pct_at_target', 0), years=p['years'])),
    ]
    for i, (label, val) in enumerate(params):
        _add_card(slide, 0.6 + i * 4.0, y, 3.7, 0.9, label, val, lang=lang)

    all_ok = vcpu_headroom >= 0 and ram_headroom >= 0 and stor_headroom >= 0
    if all_ok:
        verdict_text = t("export.pptx.projection_verdict_ok", years=p['years'])
    else:
        verdict_text = t("export.pptx.projection_verdict_expand", years=p['years'])

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.9), Inches(11.2), Inches(0.5))
    tf = txBox.text_frame
    pr = tf.paragraphs[0]
    run = pr.add_run()
    run.text = verdict_text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = GREEN if all_ok else SC_BLUE
    run.font.name = font_for(lang, run.font.name)

    if full_cluster:
        cpu_note = slide.shapes.add_textbox(Inches(0.6), Inches(6.45), Inches(11.2), Inches(0.4))
        ctf = cpu_note.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cr = cp.add_run()
        cr.text = t("export.pptx.cpu_full_cluster_note")
        cr.font.size = Pt(10.5)
        cr.font.color.rgb = CHARCOAL
        cr.font.name = font_for(lang, cr.font.name)

    disclaimer = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(11.2), Inches(0.7))
    dtf = disclaimer.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dr = dp.add_run()
    dr.text = t("export.pptx.projection_disclaimer", years=p['years'])
    dr.font.size = Pt(10)
    dr.font.italic = True
    dr.font.color.rgb = MID_GRAY
    dr.font.name = font_for(lang, dr.font.name)
